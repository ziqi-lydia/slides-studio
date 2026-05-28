#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
slides-studio: Generate high-quality, fully editable PPTX from JSON spec.
Every element (text, shape, image) is a native PPTX object — editable in
PowerPoint, Google Slides, and Keynote.

Usage:
    python create_slides.py spec.json output.pptx
"""

import json
import sys
import os
import math
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.oxml.ns import qn, nsmap
    from pptx.oxml import parse_xml
    import lxml.etree as etree
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "lxml", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.oxml.ns import qn, nsmap
    from pptx.oxml import parse_xml
    import lxml.etree as etree

# ── Helpers ──────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str):
    """Convert '#RRGGBB' or 'RRGGBB' to RGBColor."""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def parse_pos(val, default=0):
    """Parse position/size value. Supports inches (float), 'Xin', 'Xpt', 'Xemu'."""
    if val is None:
        return Inches(default)
    if isinstance(val, (int, float)):
        return Inches(val)
    s = str(val).strip().lower()
    if s.endswith('in'):
        return Inches(float(s[:-2]))
    if s.endswith('pt'):
        return Pt(float(s[:-2]))
    if s.endswith('emu'):
        return Emu(int(s[:-3]))
    return Inches(float(s))

def parse_font_size(val):
    """Parse font size to Pt."""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        return Pt(val)
    s = str(val).strip().lower()
    if s.endswith('pt'):
        return Pt(float(s[:-2]))
    return Pt(float(s))

ALIGN_MAP = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY,
}

ANCHOR_MAP = {
    'top': MSO_ANCHOR.TOP,
    'middle': MSO_ANCHOR.MIDDLE,
    'bottom': MSO_ANCHOR.BOTTOM,
}

SHAPE_MAP = {
    'rectangle': MSO_SHAPE.RECTANGLE,
    'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
    'circle': MSO_SHAPE.OVAL,
    'oval': MSO_SHAPE.OVAL,
    'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
    'diamond': MSO_SHAPE.DIAMOND,
    'pentagon': MSO_SHAPE.PENTAGON,
    'hexagon': MSO_SHAPE.HEXAGON,
    'arrow_right': MSO_SHAPE.RIGHT_ARROW,
    'arrow_left': MSO_SHAPE.LEFT_ARROW,
    'chevron': MSO_SHAPE.CHEVRON,
    'star': MSO_SHAPE.STAR_5_POINT,
    'heart': MSO_SHAPE.HEART,
    'line_horizontal': MSO_SHAPE.RECTANGLE,
    'line_vertical': MSO_SHAPE.RECTANGLE,
    'callout': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    'cloud': MSO_SHAPE.CLOUD,
    'cross': MSO_SHAPE.CROSS,
    'donut': MSO_SHAPE.DONUT,
    'no_smoking': MSO_SHAPE.NO_SYMBOL,
}

# ── Slide Background ────────────────────────────────────────────────────

def apply_background(slide, bg_spec):
    """Apply background: solid color, gradient, or image."""
    if not bg_spec:
        return

    background = slide.background
    fill = background.fill

    if 'color' in bg_spec:
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_spec['color'])

    elif 'gradient' in bg_spec:
        grad = bg_spec['gradient']
        fill.gradient()
        stops = grad.get('stops', [])
        if len(stops) >= 2:
            fill.gradient_stops[0].color.rgb = hex_to_rgb(stops[0]['color'])
            fill.gradient_stops[0].position = stops[0].get('position', 0.0)
            fill.gradient_stops[1].color.rgb = hex_to_rgb(stops[1]['color'])
            fill.gradient_stops[1].position = stops[1].get('position', 1.0)
            # Add extra stops
            for s in stops[2:]:
                gs = fill.gradient_stops.add_stop(s.get('position', 0.5))
                gs.color.rgb = hex_to_rgb(s['color'])
        angle = grad.get('angle', 0)
        fill.gradient_angle = angle

    elif 'image' in bg_spec:
        img_path = bg_spec['image']
        if os.path.exists(img_path):
            # Add image as background using XML manipulation
            from pptx.oxml.ns import qn
            bg_elem = slide.background._element
            bg_fill = bg_elem.find(qn('p:bgPr'))
            if bg_fill is None:
                bg_fill = parse_xml(
                    '<p:bgPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                    '<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                    '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                    '</p:bgPr>'
                )
                bg_elem.append(bg_fill)
            # Use slide-level image instead
            slide_width = Inches(13.333)
            slide_height = Inches(7.5)
            pic = slide.shapes.add_picture(img_path, 0, 0, slide_width, slide_height)
            # Send to back
            sp = pic._element
            sp.getparent().remove(sp)
            slide.shapes._spTree.insert(2, sp)


# ── Text Runs ───────────────────────────────────────────────────────────

def apply_run_style(run, style):
    """Apply font styling to a text run."""
    font = run.font
    if 'font_family' in style:
        font.name = style['font_family']
    if 'font_size' in style:
        font.size = parse_font_size(style['font_size'])
    if 'color' in style:
        font.color.rgb = hex_to_rgb(style['color'])
    if style.get('bold'):
        font.bold = True
    if style.get('italic'):
        font.italic = True
    if style.get('underline'):
        font.underline = True
    if 'letter_spacing' in style:
        # python-pptx doesn't directly support letter-spacing via API
        # We use XML manipulation
        spacing_val = style['letter_spacing']  # in points * 100
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(float(spacing_val) * 100)))


def add_text_to_frame(tf, text_spec):
    """Add text with rich formatting to a text frame.
    
    text_spec can be:
    - A string (simple text)
    - A dict with 'runs' (rich text with multiple styled segments)
    - A dict with 'paragraphs' (multiple paragraphs)
    """
    if isinstance(text_spec, str):
        tf.text = text_spec
        return

    if isinstance(text_spec, dict):
        paragraphs = text_spec.get('paragraphs', [])
        if not paragraphs and 'runs' in text_spec:
            paragraphs = [{'runs': text_spec['runs']}]
        if not paragraphs and 'text' in text_spec:
            tf.text = text_spec['text']
            return

        for i, para_spec in enumerate(paragraphs):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Paragraph-level formatting
            if 'align' in para_spec:
                p.alignment = ALIGN_MAP.get(para_spec['align'], PP_ALIGN.LEFT)
            if 'space_before' in para_spec:
                p.space_before = parse_font_size(para_spec['space_before'])
            if 'space_after' in para_spec:
                p.space_after = parse_font_size(para_spec['space_after'])
            if 'line_spacing' in para_spec:
                p.line_spacing = para_spec['line_spacing']

            runs = para_spec.get('runs', [])
            if isinstance(runs, str):
                runs = [{'text': runs}]

            for j, run_spec in enumerate(runs):
                if isinstance(run_spec, str):
                    run_spec = {'text': run_spec}
                if j == 0 and not p.runs:
                    run = p.add_run()
                else:
                    run = p.add_run()
                run.text = run_spec.get('text', '')
                style = {k: v for k, v in run_spec.items() if k != 'text'}
                if style:
                    apply_run_style(run, style)


# ── Element Builders ────────────────────────────────────────────────────

def add_textbox(slide, elem, design):
    """Add a text box element."""
    x = parse_pos(elem.get('x', 0))
    y = parse_pos(elem.get('y', 0))
    w = parse_pos(elem.get('w', 4))
    h = parse_pos(elem.get('h', 1))

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = elem.get('word_wrap', True)

    if elem.get('auto_size') == 'none':
        tf.auto_size = MSO_AUTO_SIZE.NONE
    elif elem.get('auto_size') == 'fit':
        tf.auto_size = MSO_AUTO_SIZE.BEST_FIT
    else:
        tf.auto_size = MSO_AUTO_SIZE.NONE

    # Margins
    if 'margin' in elem:
        m = elem['margin']
        if isinstance(m, (int, float)):
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(m)
        elif isinstance(m, dict):
            if 'left' in m: tf.margin_left = parse_pos(m['left'])
            if 'right' in m: tf.margin_right = parse_pos(m['right'])
            if 'top' in m: tf.margin_top = parse_pos(m['top'])
            if 'bottom' in m: tf.margin_bottom = parse_pos(m['bottom'])
    else:
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

    # Vertical alignment
    if 'vertical_align' in elem:
        tf.paragraphs[0]  # ensure exists
        # Set anchor
        txBody = txBox.text_frame._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            anchor_val = {'top': 't', 'middle': 'ctr', 'bottom': 'b'}.get(elem['vertical_align'], 't')
            bodyPr.set('anchor', anchor_val)

    # Add text content
    content = elem.get('content', elem.get('text', ''))
    if content:
        add_text_to_frame(tf, content)

    # Apply default styling from element level
    default_style = {}
    for key in ['font_family', 'font_size', 'color', 'bold', 'italic', 'align']:
        if key in elem:
            default_style[key] = elem[key]

    # Apply to all paragraphs/runs if not individually styled
    if default_style:
        for para in tf.paragraphs:
            if 'align' in default_style:
                para.alignment = ALIGN_MAP.get(default_style['align'], PP_ALIGN.LEFT)
            for run in para.runs:
                apply_run_style(run, default_style)

    # Rotation
    if 'rotation' in elem:
        txBox.rotation = elem['rotation']

    return txBox


def add_shape(slide, elem, design):
    """Add a shape element (rectangle, circle, etc.)."""
    shape_type = SHAPE_MAP.get(elem.get('shape', 'rectangle'), MSO_SHAPE.RECTANGLE)
    x = parse_pos(elem.get('x', 0))
    y = parse_pos(elem.get('y', 0))
    w = parse_pos(elem.get('w', 2))
    h = parse_pos(elem.get('h', 2))

    shape = slide.shapes.add_shape(shape_type, x, y, w, h)

    # Fill
    fill_spec = elem.get('fill')
    if fill_spec:
        if isinstance(fill_spec, str):
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_spec)
        elif isinstance(fill_spec, dict):
            if 'color' in fill_spec:
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(fill_spec['color'])
                if 'transparency' in fill_spec:
                    # Set alpha via XML
                    solidFill = shape.fill._fill
                    srgbClr = solidFill.find(qn('a:srgbClr'))
                    if srgbClr is not None:
                        alpha = parse_xml(
                            f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                            f'val="{int((1 - fill_spec["transparency"]) * 100000)}"/>'
                        )
                        srgbClr.append(alpha)
            elif fill_spec.get('type') == 'none':
                shape.fill.background()
    elif elem.get('no_fill'):
        shape.fill.background()

    # Line / border
    line_spec = elem.get('line', elem.get('border'))
    if line_spec:
        if isinstance(line_spec, str):
            shape.line.color.rgb = hex_to_rgb(line_spec)
            shape.line.width = Pt(1)
        elif isinstance(line_spec, dict):
            if 'color' in line_spec:
                shape.line.color.rgb = hex_to_rgb(line_spec['color'])
            if 'width' in line_spec:
                shape.line.width = parse_font_size(line_spec['width'])
            if line_spec.get('type') == 'none':
                shape.line.fill.background()
    else:
        shape.line.fill.background()

    # Shadow
    if elem.get('shadow'):
        shadow = elem['shadow']
        sp = shape._element.spPr
        effectLst = parse_xml(
            '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:outerShdw blurRad="76200" dist="38100" dir="2700000" algn="tl" rotWithShape="0">'
            f'<a:srgbClr val="{shadow.get("color", "000000")}">'
            f'<a:alpha val="{int(shadow.get("opacity", 0.3) * 100000)}"/>'
            '</a:srgbClr></a:outerShdw></a:effectLst>'
        )
        sp.append(effectLst)

    # Text in shape
    content = elem.get('content', elem.get('text'))
    if content:
        tf = shape.text_frame
        tf.word_wrap = elem.get('word_wrap', True)
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.1)
        add_text_to_frame(tf, content)
        # Apply default styling
        for key in ['font_family', 'font_size', 'color', 'bold', 'italic', 'align']:
            if key in elem:
                for para in tf.paragraphs:
                    if key == 'align':
                        para.alignment = ALIGN_MAP.get(elem[key], PP_ALIGN.LEFT)
                    for run in para.runs:
                        apply_run_style(run, {key: elem[key]})
        # Vertical alignment
        va = elem.get('vertical_align', 'middle')
        txBody = tf._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', {'top': 't', 'middle': 'ctr', 'bottom': 'b'}.get(va, 'ctr'))

    # Rounded corners
    if 'corner_radius' in elem:
        sp = shape._element.spPr
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = parse_xml('<a:avLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
                prstGeom.append(avLst)
            gd = parse_xml(
                f'<a:gd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                f'name="adj" fmla="val {int(elem["corner_radius"] * 100)}"/>'
            )
            avLst.append(gd)

    # Rotation
    if 'rotation' in elem:
        shape.rotation = elem['rotation']

    return shape


def add_image(slide, elem, design):
    """Add an image element."""
    img_path = elem.get('path', elem.get('src', elem.get('image', '')))
    if not img_path or not os.path.exists(img_path):
        print(f"  Warning: Image not found: {img_path}")
        # Add placeholder shape instead
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            parse_pos(elem.get('x', 0)),
            parse_pos(elem.get('y', 0)),
            parse_pos(elem.get('w', 3)),
            parse_pos(elem.get('h', 2)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.text = f"[Image: {os.path.basename(img_path)}]"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        return shape

    x = parse_pos(elem.get('x', 0))
    y = parse_pos(elem.get('y', 0))
    w = parse_pos(elem.get('w')) if elem.get('w') else None
    h = parse_pos(elem.get('h')) if elem.get('h') else None

    if w and h:
        pic = slide.shapes.add_picture(img_path, x, y, w, h)
    elif w:
        pic = slide.shapes.add_picture(img_path, x, y, width=w)
    elif h:
        pic = slide.shapes.add_picture(img_path, x, y, height=h)
    else:
        pic = slide.shapes.add_picture(img_path, x, y)

    # Crop
    if 'crop' in elem:
        crop = elem['crop']
        pic.crop_left = crop.get('left', 0)
        pic.crop_right = crop.get('right', 0)
        pic.crop_top = crop.get('top', 0)
        pic.crop_bottom = crop.get('bottom', 0)

    # Rotation
    if 'rotation' in elem:
        pic.rotation = elem['rotation']

    return pic


def add_line(slide, elem, design):
    """Add a line connector."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    x1 = parse_pos(elem.get('x', 0))
    y1 = parse_pos(elem.get('y', 0))
    x2 = parse_pos(elem.get('x2', elem.get('x', 0) + elem.get('w', 2)))
    y2 = parse_pos(elem.get('y2', elem.get('y', 0) + elem.get('h', 0)))

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )

    if 'color' in elem:
        connector.line.color.rgb = hex_to_rgb(elem['color'])
    if 'width' in elem:
        connector.line.width = parse_font_size(elem['width'])
    else:
        connector.line.width = Pt(1)

    return connector


def add_table(slide, elem, design):
    """Add a table element."""
    x = parse_pos(elem.get('x', 0.5))
    y = parse_pos(elem.get('y', 1.5))
    w = parse_pos(elem.get('w', 9))
    h = parse_pos(elem.get('h', 4))

    data = elem.get('data', [])
    if not data:
        return None

    rows = len(data)
    cols = max(len(row) for row in data) if data else 1

    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table

    # Style
    style = elem.get('style', {})
    header_fill = style.get('header_fill', design.get('palette', {}).get('primary', '#1E2761'))
    header_text_color = style.get('header_text_color', '#FFFFFF')
    cell_fill = style.get('cell_fill')
    alt_fill = style.get('alt_row_fill')
    border_color = style.get('border_color')
    font_family = style.get('font_family', design.get('fonts', {}).get('body', 'Calibri'))

    for i, row_data in enumerate(data):
        for j, cell_val in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_val) if cell_val is not None else ''

            # Styling
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
            else:
                run = para.add_run()
                run.text = str(cell_val) if cell_val is not None else ''
                # Clear the auto-text
                cell.text_frame.paragraphs[0].clear()
                para = cell.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = str(cell_val) if cell_val is not None else ''

            run.font.name = font_family
            run.font.size = Pt(style.get('font_size', 11))

            if i == 0 and not style.get('no_header'):
                # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(header_fill)
                run.font.color.rgb = hex_to_rgb(header_text_color)
                run.font.bold = True
            else:
                if alt_fill and i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(alt_fill)
                elif cell_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(cell_fill)

            # Cell margins
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

    # Column widths
    col_widths = elem.get('col_widths', [])
    if col_widths:
        for i, cw in enumerate(col_widths):
            if i < cols:
                table.columns[i].width = parse_pos(cw)

    return table_shape


def add_group(slide, elem, design):
    """Add grouped elements (each element added individually with offset)."""
    offset_x = parse_pos(elem.get('x', 0))
    offset_y = parse_pos(elem.get('y', 0))
    
    children = elem.get('children', [])
    shapes = []
    for child in children:
        # Adjust child positions relative to group
        if 'x' in child:
            orig_x = parse_pos(child['x'])
            child['x'] = f'{(orig_x + offset_x) / 914400}in'
        if 'y' in child:
            orig_y = parse_pos(child['y'])
            child['y'] = f'{(orig_y + offset_y) / 914400}in'
        s = add_element(slide, child, design)
        if s:
            shapes.append(s)
    return shapes


# ── Element Dispatcher ──────────────────────────────────────────────────

ELEMENT_BUILDERS = {
    'textbox': add_textbox,
    'text': add_textbox,
    'shape': add_shape,
    'rect': add_shape,
    'rectangle': add_shape,
    'circle': lambda s, e, d: add_shape(s, {**e, 'shape': 'circle'}, d),
    'oval': lambda s, e, d: add_shape(s, {**e, 'shape': 'oval'}, d),
    'image': add_image,
    'picture': add_image,
    'img': add_image,
    'line': add_line,
    'connector': add_line,
    'table': add_table,
    'group': add_group,
}

def add_element(slide, elem, design):
    """Dispatch element creation based on type."""
    elem_type = elem.get('type', 'textbox')
    builder = ELEMENT_BUILDERS.get(elem_type)
    if builder:
        return builder(slide, elem, design)
    else:
        print(f"  Warning: Unknown element type '{elem_type}', skipping")
        return None


# ── Main Build ──────────────────────────────────────────────────────────

def build_presentation(spec, output_path):
    """Build a complete PPTX from the JSON spec."""
    
    # Presentation settings
    design = spec.get('design', {})
    slides_data = spec.get('slides', [])
    
    # Create presentation
    template_path = spec.get('template')
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # Set slide dimensions (16:9 default)
    dims = spec.get('dimensions', {})
    width = dims.get('width', 13.333)
    height = dims.get('height', 7.5)
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    
    # Remove any existing slides from blank template
    if not template_path:
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(qn('r:id'))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    
    # Build each slide
    for idx, slide_spec in enumerate(slides_data):
        print(f"  Building slide {idx + 1}: {slide_spec.get('title', 'Untitled')}")
        
        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        bg = slide_spec.get('background', design.get('default_background'))
        if bg:
            apply_background(slide, bg)
        
        # Elements
        elements = slide_spec.get('elements', [])
        for elem in elements:
            add_element(slide, elem, design)
        
        # Speaker notes
        notes = slide_spec.get('notes', '')
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
    
    # Save
    prs.save(output_path)
    print(f"\n✅ Saved: {output_path} ({len(slides_data)} slides)")
    return output_path


# ── CLI ─────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("Usage: python create_slides.py <spec.json> <output.pptx>")
        sys.exit(1)

    spec_path = sys.argv[1]
    output_path = sys.argv[2]

    with open(spec_path, 'r', encoding='utf-8') as f:
        spec = json.load(f)

    build_presentation(spec, output_path)


if __name__ == "__main__":
    main()
